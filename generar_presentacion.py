"""
Genera la presentación V1 del Chatbot de Autodiagnóstico (.pptx).

No forma parte del chatbot: es una utilidad para producir el archivo de
presentación a partir de los datos y el estado real del proyecto.

Uso:
    .venv\\Scripts\\python generar_presentacion.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# --- Paleta y tipografía ------------------------------------------------------
INK = RGBColor(0x0E, 0x24, 0x30)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
AMBER = RGBColor(0xF2, 0xA3, 0x40)
BONE = RGBColor(0xF6, 0xF3, 0xEC)
SAGE = RGBColor(0x5C, 0x7A, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK_SOFT = RGBColor(0x22, 0x3B, 0x46)
CARD_BG = RGBColor(0xEC, 0xE8, 0xDD)

F_HEAD = "Century Gothic"
F_BODY = "Calibri"
F_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _no_autofit(tf):
    """Evita que PowerPoint reduzca el texto solo (mantenemos control del tamaño)."""
    from lxml import etree
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        n = bodyPr.find(qn(tag))
        if n is not None:
            bodyPr.remove(n)
    if bodyPr.find(qn('a:noAutofit')) is None:
        etree.SubElement(bodyPr, qn('a:noAutofit'))


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=INK, font=F_BODY,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, letter_spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
        r.font.bold = bold
        r.font.italic = italic
    return box


def add_eyebrow_rail(slide, num_text, section_text, dark=True):
    """Franja lateral de acento con número de sección — marca la navegación."""
    bg = INK if dark else BONE
    fg = TEAL if dark else INK
    sub = RGBColor(0x9F, 0xB4, 0xB8) if dark else SAGE
    add_rect(slide, 0, 0, Inches(1.15), SLIDE_H, bg)
    add_text(slide, Inches(0.15), Inches(0.55), Inches(0.9), Inches(0.6),
              num_text, size=30, color=fg, font=F_HEAD, bold=True)
    add_text(slide, Inches(0.15), Inches(6.5), Inches(0.9), Inches(0.8),
              section_text, size=9, color=sub, font=F_BODY,
              align=PP_ALIGN.LEFT)


def add_title(slide, kicker, title, x=Inches(1.55), y=Inches(0.55), w=Inches(11.4)):
    add_text(slide, x, y, w, Inches(0.4), kicker.upper(), size=13, color=TEAL,
              font=F_BODY, bold=True)
    add_text(slide, x, y + Inches(0.38), w, Inches(0.95), title, size=34,
              color=INK, font=F_HEAD, bold=True, line_spacing=1.0)


def add_card(slide, x, y, w, h, title, body_lines, title_color=INK, accent=TEAL):
    add_rect(slide, x, y, w, h, CARD_BG)
    add_rect(slide, x, y, Inches(0.06), h, accent)
    add_text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5), Inches(0.4),
              title, size=15, color=title_color, font=F_HEAD, bold=True)
    body = "\n".join(f"•  {b}" for b in body_lines)
    add_text(slide, x + Inches(0.28), y + Inches(0.62), w - Inches(0.5), h - Inches(0.8),
              body, size=11.5, color=INK_SOFT, font=F_BODY, line_spacing=1.15)


def add_chip(slide, x, y, w, h, text, bg=WHITE, fg=INK, border=None, size=11.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    _no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = fg
    r.font.name = F_BODY
    r.font.bold = True
    return shp


def add_footer(slide, page_no):
    add_text(slide, Inches(11.8), Inches(7.1), Inches(1.3), Inches(0.3),
              f"{page_no:02d}", size=10, color=SAGE, font=F_MONO,
              align=PP_ALIGN.RIGHT)


# ==============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- Slide 1: Portada
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)
add_rect(s, 0, Inches(6.55), SLIDE_W, Inches(0.05), TEAL)
add_text(s, Inches(1.1), Inches(2.55), Inches(9.5), Inches(0.5),
          "AUTODIAGNÓSTICO · CHATBOT DE IA", size=15, color=TEAL, font=F_BODY, bold=True)
add_text(s, Inches(1.05), Inches(3.0), Inches(11), Inches(1.7),
          "Versión 1", size=64, color=WHITE, font=F_HEAD, bold=True)
add_text(s, Inches(1.1), Inches(4.55), Inches(10.5), Inches(0.6),
          "Chatbot en lenguaje natural conectado a datos reales del proceso de autodiagnóstico",
          size=16, color=RGBColor(0xC9, 0xD8, 0xDA), font=F_BODY)
add_text(s, Inches(1.1), Inches(6.75), Inches(6), Inches(0.4),
          "Julio 2026", size=11, color=SAGE, font=F_MONO)
# marca gráfica: barras de "señal"
bars_x = Inches(10.6)
for i, h in enumerate([0.35, 0.6, 0.9, 1.25, 0.75]):
    add_rect(s, bars_x + Inches(i * 0.45), Inches(2.0 + (1.25 - h)), Inches(0.28),
             Inches(h), TEAL if i != 3 else AMBER)

# ---------------------------------------------------------------- Slide 2: Contexto
s = prs.slides.add_slide(BLANK)
add_bg(s, BONE)
add_eyebrow_rail(s, "01", "CONTEXTO", dark=True)
add_title(s, "Por qué existe", "Del dato disperso a la respuesta en segundos")
add_text(s, Inches(1.55), Inches(1.85), Inches(5.6), Inches(4.6),
          "El proceso de autodiagnóstico corre en tres canales distintos "
          "(Portal, WhatsApp y Sysbrazo) y su resultado —si falla, escala o "
          "se resuelve— vive repartido entre Sysbrazo y Odoo.\n\n"
          "Hoy, responder “¿por qué falla más el Portal?” o "
          "“¿qué equipo resuelve más rápido?” exige escribir SQL.\n\n"
          "El chatbot elimina esa barrera: cualquier persona pregunta con sus "
          "propias palabras y obtiene la cifra, la tabla y el gráfico al instante.",
          size=15, color=INK_SOFT, font=F_BODY, line_spacing=1.3)
add_card(s, Inches(7.5), Inches(1.85), Inches(5.3), Inches(4.6),
         "En una frase",
         ["Convierte la data del autodiagnóstico en respuestas",
          "que cualquiera puede obtener preguntando",
          "con sus propias palabras."],
         accent=TEAL)
add_footer(s, 2)

# ---------------------------------------------------------------- Slide 3: Fuente de datos (overview)
s = prs.slides.add_slide(BLANK)
add_bg(s, BONE)
add_eyebrow_rail(s, "02", "FUENTE DE DATOS", dark=True)
add_title(s, "Composición de la fuente", "Una fila por autodiagnóstico, 26 campos")
add_text(s, Inches(1.55), Inches(1.85), Inches(11.2), Inches(0.5),
          "Consulta consolidada en Redash sobre analytics.v_auto_diagnostic_full "
          "· ~72.800 registros · desde nov-2025",
          size=13, color=SAGE, font=F_BODY, italic=True)

cards_y = Inches(2.55)
card_w = Inches(3.55)
gap = Inches(0.2)
add_card(s, Inches(1.55), cards_y, card_w, Inches(3.9),
         "① Cliente y ubicación",
         ["id, client_id, client_name",
          "gaiia_id (id en sistemas internos)",
          "nombre_ciudad (6 ciudades, Colombia)"],
         accent=TEAL)
add_card(s, Inches(1.55) + card_w + gap, cards_y, card_w, Inches(3.9),
         "② El proceso en sí",
         ["source (canal): Portal, WhatsApp, Sysbrazo",
          "status (resultado): finished / failed / canceled",
          "started_at, finished_at, duración",
          "failure_reason, failed_step (causa técnica)"],
         accent=AMBER)
add_card(s, Inches(1.55) + 2 * (card_w + gap), cards_y, card_w, Inches(3.9),
         "③ El ticket (si escaló)",
         ["odoo_ticket_id, ticket_ref, ticket_name",
          "ticket_stage: New / In Progress / Solved",
          "ticket_team: NOC, Instalaciones y Mtto., CX...",
          "motivo apertura / cierre",
          "fechas de apertura y cierre, horas de resolución"],
         accent=INK)
add_footer(s, 3)

# ---------------------------------------------------------------- Slide 4: Fuente de datos (proporciones reales)
s = prs.slides.add_slide(BLANK)
add_bg(s, BONE)
add_eyebrow_rail(s, "02", "FUENTE DE DATOS", dark=True)
add_title(s, "Composición de la fuente", "Lo que hay realmente en los datos hoy")

# Tabla-resumen tipo "stat"
stats = [
    ("72.827", "Autodiagnósticos"),
    ("3", "Canales"),
    ("6", "Ciudades (Colombia)"),
    ("9.742", "Con ticket asociado"),
    ("6.984", "Tickets ya resueltos"),
]
sx = Inches(1.55)
sw = Inches(2.15)
for i, (val, lbl) in enumerate(stats):
    x = sx + i * (sw + Inches(0.08))
    add_rect(s, x, Inches(1.9), sw, Inches(1.35), INK)
    add_text(s, x, Inches(2.0), sw, Inches(0.7), val, size=26, color=TEAL,
              font=F_MONO, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(2.68), sw - Inches(0.2), Inches(0.5), lbl,
              size=10.5, color=WHITE, font=F_BODY, align=PP_ALIGN.CENTER)

add_text(s, Inches(1.55), Inches(3.6), Inches(5.5), Inches(0.4),
          "CANALES (source)", size=12, color=SAGE, font=F_BODY, bold=True)
canales = [("Portal web", "≈ 73k"), ("WhatsApp (Botmaker)", "≈ 4,6k"), ("Sysbrazo", "≈ 0,6k")]
cy = Inches(4.05)
for nombre, val in canales:
    add_rect(s, Inches(1.55), cy, Inches(5.5), Inches(0.55), CARD_BG)
    add_text(s, Inches(1.75), cy + Inches(0.08), Inches(3.2), Inches(0.4), nombre,
              size=13, color=INK, font=F_BODY, bold=True)
    add_text(s, Inches(5.3), cy + Inches(0.08), Inches(1.6), Inches(0.4), val,
              size=13, color=TEAL, font=F_MONO, bold=True, align=PP_ALIGN.RIGHT)
    cy += Inches(0.68)

add_text(s, Inches(7.5), Inches(3.6), Inches(5.3), Inches(0.4),
          "CIUDADES CUBIERTAS", size=12, color=SAGE, font=F_BODY, bold=True)
ciudades = ["Barranquilla", "Cartagena", "Montería", "Santa Marta", "Sincelejo", "Turbaco"]
cw = Inches(1.68)
for i, ciu in enumerate(ciudades):
    col = i % 3
    row = i // 3
    x = Inches(7.5) + col * (cw + Inches(0.08))
    y = Inches(4.05) + row * Inches(0.62)
    add_chip(s, x, y, cw, Inches(0.5), ciu, bg=WHITE, fg=INK, border=TEAL, size=11)

add_text(s, Inches(7.5), Inches(5.55), Inches(5.3), Inches(1.2),
          "Nota: se excluyó “Corrientes” (fuera del alcance geográfico, "
          "análisis limitado a Colombia).",
          size=11, color=SAGE, font=F_BODY, italic=True, line_spacing=1.2)
add_footer(s, 4)

# ---------------------------------------------------------------- Slide 5: Alcance V1
s = prs.slides.add_slide(BLANK)
add_bg(s, BONE)
add_eyebrow_rail(s, "03", "ALCANCE V1", dark=True)
add_title(s, "Qué puede y qué no puede (todavía)", "Alcance de esta primera versión")

add_rect(s, Inches(1.55), Inches(1.85), Inches(5.5), Inches(0.55), TEAL)
add_text(s, Inches(1.75), Inches(1.93), Inches(5), Inches(0.4), "✓  INCLUIDO EN V1",
          size=14, color=WHITE, font=F_HEAD, bold=True)
incluido = [
    "Preguntas en lenguaje natural, sin sintaxis fija",
    "Memoria conversacional (sigue el hilo de la charla)",
    "3 canales · 6 ciudades · histórico desde nov-2025",
    "Resultado del proceso + causa técnica de falla",
    "Ticket asociado: equipo, estado y tiempo de resolución",
    "Filtros y gráficos automáticos por pregunta",
]
add_text(s, Inches(1.75), Inches(2.55), Inches(5.2), Inches(4),
          "\n".join(f"•  {t}" for t in incluido), size=13.5, color=INK_SOFT,
          font=F_BODY, line_spacing=1.35)

add_rect(s, Inches(7.5), Inches(1.85), Inches(5.3), Inches(0.55), AMBER)
add_text(s, Inches(7.7), Inches(1.93), Inches(5), Inches(0.4), "◐  FUERA DE ALCANCE (AÚN)",
          size=14, color=INK, font=F_HEAD, bold=True)
fuera = [
    "El comentario/mensaje exacto entregado al cliente",
    "Actualización 100% automática (hoy es manual, con botón)",
    "Canales fuera de esta fuente (si los hubiera)",
    "Acceso privado activado por defecto (existe, se activa a pedido)",
]
add_text(s, Inches(7.7), Inches(2.55), Inches(4.9), Inches(4),
          "\n".join(f"•  {t}" for t in fuera), size=13.5, color=INK_SOFT,
          font=F_BODY, line_spacing=1.35)
add_footer(s, 5)

# ---------------------------------------------------------------- Slide 6: Filtros y botones
s = prs.slides.add_slide(BLANK)
add_bg(s, BONE)
add_eyebrow_rail(s, "04", "USO", dark=True)
add_title(s, "La barra de control", "Filtros y botones disponibles")

add_text(s, Inches(1.55), Inches(1.85), Inches(5.5), Inches(0.4),
          "FILTROS (barra izquierda)", size=13, color=SAGE, font=F_BODY, bold=True)
filtros = [
    ("Estado del ticket", "New · In Progress · Solved"),
    ("Canal", "Portal · WhatsApp · Sysbrazo"),
    ("Ciudad", "las 6 ciudades disponibles"),
    ("Fecha inicio / Hasta", "cualquier rango dentro del histórico"),
]
fy = Inches(2.35)
for nombre, detalle in filtros:
    add_rect(s, Inches(1.55), fy, Inches(5.6), Inches(0.85), CARD_BG)
    add_rect(s, Inches(1.55), fy, Inches(0.06), Inches(0.85), TEAL)
    add_text(s, Inches(1.8), fy + Inches(0.1), Inches(5.1), Inches(0.35), nombre,
              size=14, color=INK, font=F_BODY, bold=True)
    add_text(s, Inches(1.8), fy + Inches(0.45), Inches(5.2), Inches(0.35), detalle,
              size=11.5, color=SAGE, font=F_BODY)
    fy += Inches(1.0)

add_text(s, Inches(7.5), Inches(1.85), Inches(5.3), Inches(0.4),
          "BOTONES", size=13, color=SAGE, font=F_BODY, bold=True)
botones = [
    ("🔄  Actualizar datos", "Trae el último resultado ya calculado en Redash"),
    ("🔄  Reiniciar chat", "Borra la conversación y la memoria, sin recargar"),
    ("🔒  Acceso (opcional)", "Pantalla de usuario/clave antes de entrar, si se activa"),
]
by = Inches(2.35)
for nombre, detalle in botones:
    add_rect(s, Inches(7.5), by, Inches(5.3), Inches(1.05), INK)
    add_text(s, Inches(7.75), by + Inches(0.12), Inches(4.9), Inches(0.4), nombre,
              size=15, color=TEAL, font=F_BODY, bold=True)
    add_text(s, Inches(7.75), by + Inches(0.55), Inches(4.9), Inches(0.45), detalle,
              size=11.5, color=RGBColor(0xC9, 0xD8, 0xDA), font=F_BODY, line_spacing=1.1)
    by += Inches(1.2)
add_footer(s, 6)

# ---------------------------------------------------------------- Slide 7: Arquitectura técnica
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)
add_eyebrow_rail(s, "05", "POR DENTRO", dark=True)
add_title(s, "Explicación técnica (mini)", "Cómo viaja la pregunta hasta la respuesta",
          x=Inches(1.55))
# recolor title text since bg is dark
for shp in list(s.shapes)[-2:]:
    pass  # (kept simple; título ya usa TEAL/INK, ajustamos abajo con cajas propias)

# Reemplazo simple del título para fondo oscuro (blanco) - se sobreescribe visualmente
add_rect(s, Inches(1.55), Inches(0.55), Inches(11.4), Inches(1.3), INK)  # tapa el título por defecto (fondo oscuro)
add_text(s, Inches(1.55), Inches(0.55), Inches(11.4), Inches(0.4),
          "POR DENTRO", size=13, color=TEAL, font=F_BODY, bold=True)
add_text(s, Inches(1.55), Inches(0.93), Inches(11.4), Inches(0.85),
          "Cómo viaja la pregunta hasta la respuesta", size=32, color=WHITE,
          font=F_HEAD, bold=True)

pasos = [
    ("1", "Redash", "Consulta SQL ya\nconstruida y validada"),
    ("2", "API", "El chatbot trae el\nresultado vía API"),
    ("3", "Chatbot", "Carga los datos y\nespera tu pregunta"),
    ("4", "IA (Gemini)", "Traduce la pregunta\na una consulta"),
    ("5", "Candado", "Se valida: solo\nlectura, nunca escribe"),
    ("6", "Respuesta", "Cifra + tabla +\ngráfico, en español"),
]
box_w = Inches(1.75)
box_h = Inches(2.0)
total_w = box_w * 6 + Inches(0.25) * 5
start_x = (SLIDE_W - total_w) / 2 + Inches(0.3)
y0 = Inches(2.6)
for i, (n, titulo, desc) in enumerate(pasos):
    x = start_x + i * (box_w + Inches(0.25))
    color = TEAL if i % 2 == 0 else AMBER
    add_rect(s, x, y0, box_w, box_h, INK_SOFT)
    add_rect(s, x, y0, box_w, Inches(0.06), color)
    add_text(s, x + Inches(0.12), y0 + Inches(0.18), box_w - Inches(0.24), Inches(0.4),
              n, size=20, color=color, font=F_MONO, bold=True)
    add_text(s, x + Inches(0.12), y0 + Inches(0.62), box_w - Inches(0.24), Inches(0.4),
              titulo, size=13.5, color=WHITE, font=F_HEAD, bold=True)
    add_text(s, x + Inches(0.12), y0 + Inches(1.05), box_w - Inches(0.24), Inches(0.9),
              desc, size=10, color=RGBColor(0xB9, 0xCC, 0xCF), font=F_BODY, line_spacing=1.15)
    if i < 5:
        arrow_x = x + box_w
        add_text(s, arrow_x, y0 + Inches(0.7), Inches(0.25), Inches(0.6), "→",
                  size=20, color=TEAL, align=PP_ALIGN.CENTER)

add_text(s, Inches(1.55), Inches(5.1), Inches(10.5), Inches(0.5),
          "El SQL que ejecutó cada respuesta siempre queda visible "
          "(“¿Cómo lo calculé?”) — las cifras nunca se inventan.",
          size=13, color=RGBColor(0xC9, 0xD8, 0xDA), font=F_BODY, italic=True)

add_text(s, Inches(1.55), Inches(5.75), Inches(10.7), Inches(1.3),
          "Fuente de datos: Redash (consulta sobre analytics.v_auto_diagnostic_full)   ·   "
          "Motor de análisis: DuckDB (en memoria, de solo lectura)   ·   "
          "Interfaz: Streamlit (Python)   ·   Modelo de IA: Google Gemini",
          size=11, color=SAGE, font=F_MONO, line_spacing=1.3)
add_footer(s, 7)

# ---------------------------------------------------------------- Slide 8: Qué sigue (V2)
s = prs.slides.add_slide(BLANK)
add_bg(s, BONE)
add_eyebrow_rail(s, "06", "PRÓXIMOS PASOS", dark=True)
add_title(s, "Qué sigue", "Hoja de ruta hacia la V2")

# Bloque destacado: el pendiente principal
add_rect(s, Inches(1.55), Inches(1.85), Inches(11.25), Inches(2.15), INK)
add_rect(s, Inches(1.55), Inches(1.85), Inches(0.08), Inches(2.15), AMBER)
add_text(s, Inches(1.85), Inches(2.02), Inches(4), Inches(0.4),
          "PRIORIDAD #1", size=12, color=AMBER, font=F_BODY, bold=True)
add_text(s, Inches(1.85), Inches(2.4), Inches(10.7), Inches(0.55),
          "Trazar el mensaje final entregado al cliente", size=22, color=WHITE,
          font=F_HEAD, bold=True)
add_text(s, Inches(1.85), Inches(2.95), Inches(10.7), Inches(0.95),
          "Hoy el sistema NO guarda el texto exacto que recibe el cliente tras un "
          "autodiagnóstico (éxito o falla): se arma en el momento a partir de "
          "plantillas o se delega a Botmaker, y ese texto resuelto no se persiste "
          "en ninguna tabla. La V2 debe definir dónde y cómo capturarlo.",
          size=13, color=RGBColor(0xD8, 0xE4, 0xE6), font=F_BODY, line_spacing=1.25)

add_text(s, Inches(1.55), Inches(4.3), Inches(6), Inches(0.4),
          "TAMBIÉN EN EL RADAR", size=13, color=SAGE, font=F_BODY, bold=True)
radar = [
    "Refresco automático programado (sin depender del botón)",
    "Confirmar si “NET Operations” y “NOC” son el mismo equipo",
    "Detalle de motivo de cierre por ticket",
    "Activar acceso privado en producción",
]
add_text(s, Inches(1.55), Inches(4.85), Inches(11), Inches(2.2),
          "\n".join(f"•  {t}" for t in radar), size=13.5, color=INK_SOFT,
          font=F_BODY, line_spacing=1.4)
add_footer(s, 8)

# ---------------------------------------------------------------- Slide 9: Cierre
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)
add_rect(s, 0, Inches(3.6), SLIDE_W, Inches(0.05), TEAL)
add_text(s, Inches(1.1), Inches(2.7), Inches(10), Inches(0.9),
          "Gracias", size=48, color=WHITE, font=F_HEAD, bold=True)
add_text(s, Inches(1.1), Inches(3.8), Inches(10), Inches(0.5),
          "Preguntas y siguiente ciclo de pruebas", size=16,
          color=RGBColor(0xC9, 0xD8, 0xDA), font=F_BODY)

prs.save("Chatbot_Autodiagnostico_V1.pptx")
print("Presentación generada: Chatbot_Autodiagnostico_V1.pptx")
